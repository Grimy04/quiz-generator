import os
import json
import random
import time
import uuid
import warnings
import requests
import re
import sys
import logging
from datetime import datetime

from cryptography.utils import CryptographyDeprecationWarning
warnings.filterwarnings("ignore", category=CryptographyDeprecationWarning)
warnings.filterwarnings("ignore", category=UserWarning, module="pypdf")

# Zittisce i noiosi log di pypdf (come "Ignoring wrong pointing object")
logging.getLogger("pypdf").setLevel(logging.ERROR)

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from google import genai
from google.genai import types
from groq import Groq
from dotenv import load_dotenv

# =========================
# SETUP API E MODELLI
# =========================

load_dotenv()

api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    print("❌ ERRORE: GEMINI_API_KEY mancante nel file .env")
    exit()

client = genai.Client(api_key=api_key)

groq_api_key = os.getenv("GROQ_API_KEY")
groq_client = Groq(api_key=groq_api_key) if groq_api_key else None

openrouter_api_key = os.getenv("OPENROUTER_API_KEY")

OPENROUTER_MODELS = [
    "nvidia/nemotron-3-ultra-550b-a55b:free",
    "nex-agi/nex-n2-pro:free",
    "qwen/qwen3-next-80b-a3b-instruct:free",
    "openai/gpt-oss-120b:free"
]

PALETTE_OPZIONI = [
    {"nome": "Amber Warm", "ground": "#100E0A", "ground-2": "#1A1712", "cream": "#ECE5D6", "cream-2": "#A69C87", "cream-3": "#766C59", "signal": "#E0913A", "signal-2": "#F4B662", "signal-ink": "#B16C26", "cap_tones": ["#DDBE73", "#CBA158", "#B9853F", "#A56A33", "#8C552B", "#9A6A2F"]},
    {"nome": "Ocean Blue", "ground": "#0A0D10", "ground-2": "#12161A", "cream": "#E0E7F3", "cream-2": "#8797A6", "cream-3": "#596876", "signal": "#3A82E0", "signal-2": "#62A4F4", "signal-ink": "#265CB1", "cap_tones": ["#739DDD", "#5887CB", "#3F6FB9", "#335AA5", "#2B4C8C", "#2F509A"]},
    {"nome": "Sage Green", "ground": "#0A100B", "ground-2": "#121A13", "cream": "#E0F3E3", "cream-2": "#87A68D", "cream-3": "#59765F", "signal": "#42704C", "signal-2": "#6BA877", "signal-ink": "#2B5234", "cap_tones": ["#73DD84", "#58CB6C", "#3FB955", "#33A548", "#2B8C3E", "#2F9A43"]}
]

# =========================
# RATE LIMIT E TEMPI
# =========================

MAX_RPM = 10
INTERVALLO = 60 / MAX_RPM
ultima_richiesta = 0
gemini_resume_time = 0

def rate_limit():
    global ultima_richiesta
    delta = time.time() - ultima_richiesta
    if delta < INTERVALLO: time.sleep(INTERVALLO - delta)
    ultima_richiesta = time.time()

# =========================
# UI: DASHBOARD RENDERER E MENU
# =========================

def print_banner():
    print(f"\033[96m╔════════════════════════════════════════════════════╗\033[0m")
    print(f"\033[96m║              AI QUIZ GENERATOR v2.0                ║\033[0m")
    print(f"\033[96m║      PDF • DOCX • PPTX → Quiz Interattivi          ║\033[0m")
    print(f"\033[96m╚════════════════════════════════════════════════════╝\033[0m\n")

def scegli_grafica():
    print("\033[93m🎨 SCEGLI L'ANTEPRIMA DELLO STILE GRAFICO\033[0m")
    print("--------------------------------------------------")
    print(" \033[97m[1] Originale 3D\033[0m   -> Stile elegante, carte inclinate con effetto profondità.")
    print(" \033[97m[2] Neumorfismo\033[0m    -> Stile app moderna, colori scuri vibranti, pannelli in vetro.")
    print(" \033[97m[3] Hacker Terminal\033[0m-> Stile Matrix, schermo CRT con sfarfallio, linea di comando.")
    print(" \033[97m[4] Hacker Terminal\033[0m-> Stile Matrix, schermo CRT con sfarfallio, linea di comando.")
    print("--------------------------------------------------")
    
    while True:
        scelta = input("👉 Digita il numero (1, 2 o 3) e premi Invio: ").strip()
        if scelta == "1":
            return "index_1.html"
        elif scelta == "2":
            return "index_2.html"
        elif scelta == "3":
            return "index_3.html"
        elif scelta == "4":
            return "index_4.html"
        else:
            print("\033[91m⚠️  Scelta non valida. Inserisci 1, 2 o 3.\033[0m")

def print_final_box(mins, files, blocks, qs, tot_db):
    def rpad(s, w=27): return str(s).ljust(w)
    tot_db_str = f"{tot_db:,}".replace(',', '.')
    
    print(f"\n\033[92m╔══════════════════════════════════════════════╗\033[0m")
    print(f"\033[92m║                ELABORAZIONE TERMINATA        ║\033[0m")
    print(f"\033[92m╠══════════════════════════════════════════════╣\033[0m")
    print(f"\033[92m║ Tempo Totale      {rpad(f'{mins} min')}║\033[0m")
    print(f"\033[92m║ File Elaborati    {rpad(files)}║\033[0m")
    print(f"\033[92m║ Blocchi           {rpad(blocks)}║\033[0m")
    print(f"\033[92m║ Domande Generate  {rpad(qs)}║\033[0m")
    print(f"\033[92m║ Database Totale   {rpad(tot_db_str)}║\033[0m")
    print(f"\033[92m╚══════════════════════════════════════════════╝\033[0m\n")

def draw_dashboard(doc_name, blocco_curr, blocco_tot, dom_gen, motore, eta, file_curr, file_tot, b_curr, b_tot, tot_dom, speed, api_gem, api_groq, api_or, ultimo_errore, is_first=False):
    if not is_first:
        sys.stdout.write("\033[18A") 

    doc_str = f"{doc_name[:28]}"
    prog_str = f"{blocco_curr} / {blocco_tot} blocchi"
    dom_str = f"{dom_gen} generate"
    orario = datetime.now().strftime("%H:%M:%S")

    print(f"┌──────────────────────────────────────────────┐\033[K")
    print(f"│ Documento     {doc_str:<30} │\033[K")
    print(f"│ Progresso     {prog_str:<30} │\033[K")
    print(f"│ Domande       {dom_str:<30} │\033[K")
    print(f"│ Motore        {motore:<30} │\033[K")
    print(f"│ ETA           {eta:<30} │\033[K")
    print(f"│ Aggiornato    {orario:<30} │\033[K")
    print(f"└──────────────────────────────────────────────┘\033[K")
    print(f"━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\033[K")
    print(f"📚 File analizzati      {file_curr} / {file_tot}\033[K")
    print(f"🧩 Blocchi processati   {b_curr} / {b_tot}\033[K")
    print(f"📝 Domande generate     {tot_dom}\033[K")
    print(f"⚡ Velocità media       {speed} blocchi/min\033[K")
    print(f"🤖 Gemini               {api_gem} richieste\033[K")
    print(f"🤖 Groq                 {api_groq} richieste\033[K")
    print(f"🤖 OpenRouter           {api_or} richieste\033[K")
    print(f"━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\033[K")
    
    err_str = ultimo_errore[:45].replace('\n', ' ') if ultimo_errore else "In ascolto... OK"
    colore_err = "\033[91m" if ultimo_errore else "\033[92m"
    print(f"{colore_err}⚠️ Status API: {err_str:<32}\033[0m\033[K")
    
    sys.stdout.flush()

# =========================
# PROMPT E PARSING
# =========================

def get_prompt(blocco):
    return f"""
        Sei un professore universitario esigente. Il tuo obiettivo è valutare la reale comprensione analitica degli studenti.
        Genera da 4 a 7 domande FONDAMENTALI basandoti ESCLUSIVAMENTE sul testo fornito.

        REGOLE DI CONTENUTO:
        1. IGNORA TOTALMENTE info organizzative: nomi prof, università, bibliografia, indici.
        2. Concentrati sui concetti teorici, tecnici e pratici. Cerca di valutare il ragionamento, non solo la memoria meccanica.
        3. ASSOLUTAMENTE VIETATO creare domande che facciano riferimento a immagini, grafici, tabelle o figure.
        4. NON ripetere mai lo stesso concetto in domande diverse.
        5. FORMATTAZIONE (CRITICO): Usa SEMPRE il Markdown per la formattazione dei campi di testo (q, opts, explain, points).
           - Usa i `backtick` per parole chiave tecniche.
           - Usa i blocchi di codice ```linguaggio ... ``` per snippet a più righe.
           - Usa il LaTeX matematico racchiuso tra $ per l'inline (es. $x^2$) e $$ per i blocchi matematici su più righe.
           - Usa il **grassetto** per evidenziare i concetti cruciali.

        REGOLE PER I FORMATI DELLE DOMANDE:
        - VERO/FALSO ("vf"): DEVONO essere affermazioni dichiarative. ASSOLUTAMENTE VIETATO usare il punto interrogativo.
        - SCELTA MULTIPLA ("mc"): Le opzioni errate DEVONO essere plausibili e ingannevoli.
        - APERTE ("open"): Richiedi di spiegare il "perché" o il "come" di un concetto chiave.

        REGOLE DI OUTPUT:
        Restituisci ESCLUSIVAMENTE un oggetto JSON valido con UNA SOLA CHIAVE chiamata "domande".
        NON inserire commenti nel JSON. Usa i blocchi markdown per racchiudere il JSON.

        Esempio di struttura JSON attesa:
        {{
        "domande": [
        {{
            "q": "Testo della domanda",
            "fmt": "vf",
            "correct": true,
            "opts": null,
            "explain": "Spiegazione della risposta",
            "points": null
        }}
        ]
        }}

        TESTO DEL BLOCCO:
        {blocco}
    """

def parse_domande(data, cid):
    items = data.get("domande", data) if isinstance(data, dict) else data
    out = []
    if not isinstance(items, list): return []
    for d in items:
        if isinstance(d, dict) and "q" in d:
            d["ch"] = cid
            d["id"] = f"{cid}-{uuid.uuid4().hex[:8]}"
            out.append(d)
    return out

# =========================
# LETTURA, PULIZIA E SPLIT LOGICO
# =========================

def pulisci_testo(testo):
    testo = re.sub(r'\n{3,}', '\n\n', testo)
    testo = re.sub(r' {2,}', ' ', testo)
    testo = re.sub(r'(?m)^\s*\d+\s*$', '', testo)
    return testo.strip()

def estrai_testo(filepath):
    ext = filepath.lower().split(".")[-1]
    testo = ""
    try:
        if ext == "pdf":
            reader = PdfReader(filepath)
            for p in reader.pages:
                t = p.extract_text()
                if t: testo += t + "\n"
        elif ext == "docx":
            doc = Document(filepath)
            for p in doc.paragraphs: testo += p.text + "\n"
        elif ext == "pptx":
            prs = Presentation(filepath)
            for s in prs.slides:
                for shape in s.shapes:
                    if hasattr(shape, "text"): testo += shape.text + "\n"
    except: pass
    
    return pulisci_testo(testo)

def split_logico(testo, max_chars=7500):
    paragrafi = testo.replace("\r\n", "\n").split("\n\n")
    blocchi = []
    blocco_corrente = ""
    for p in paragrafi:
        if len(p) > max_chars:
            if blocco_corrente.strip():
                blocchi.append(blocco_corrente.strip())
                blocco_corrente = ""
            blocchi.append(p[:max_chars].strip())
            continue
        if len(blocco_corrente) + len(p) < max_chars:
            blocco_corrente += p + "\n\n"
        else:
            if blocco_corrente.strip():
                blocchi.append(blocco_corrente.strip())
            blocco_corrente = p + "\n\n"
    if blocco_corrente.strip():
        blocchi.append(blocco_corrente.strip())
    return blocchi

# =========================
# MOTORI AI CON GESTIONE ERRORI E ROTAZIONE
# =========================

def genera_gemini(blocco, cid):
    global gemini_resume_time
    if time.time() < gemini_resume_time: return None, "Gemini", "In pausa (Rate Limit)"
    try:
        rate_limit()
        r = client.models.generate_content(
            model="gemini-2.5-flash", contents=get_prompt(blocco),
            config=types.GenerateContentConfig(response_mime_type="application/json")
        )
        return parse_domande(json.loads(r.text), cid), "Gemini", None
    except Exception as e:
        err_str = str(e)
        if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
            gemini_resume_time = time.time() + 86400
            return None, "Gemini", "Quota Esaurita (429)"
        return None, "Gemini", f"JSON Err: {err_str[:30]}"

def genera_groq(blocco, cid):
    if not groq_client: return None, "Groq", "API Key mancante"
    try:
        r = groq_client.chat.completions.create(
            messages=[{"role":"user","content":get_prompt(blocco)}],
            model="llama-3.3-70b-versatile", response_format={"type":"json_object"}
        )
        return parse_domande(json.loads(r.choices[0].message.content), cid), "Groq", None
    except Exception as e: 
        return None, "Groq", f"Err: {str(e)[:30]}"

def genera_openrouter(blocco, cid):
    if not openrouter_api_key: return None, "OpenRouter", "API Key mancante"
    ultimo_errore = ""
    for model in OPENROUTER_MODELS:
        try:
            r = requests.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {openrouter_api_key}"},
                json={"model": model, "messages":[{"role":"user","content":get_prompt(blocco)}]},
                timeout=30
            )
            if r.status_code == 200:
                content = r.json()["choices"][0]["message"]["content"]
                match = re.search(r'\{[\s\S]*\}', content)
                content_pulito = match.group(0) if match else content.replace("```json", "").replace("```", "").strip()
                return parse_domande(json.loads(content_pulito), cid), "OpenRouter", None
            elif r.status_code == 429:
                ultimo_errore = f"429 su {model.split('/')[1][:10]}"
                continue 
            else:
                ultimo_errore = f"HTTP {r.status_code}"
                continue
        except Exception as e: 
            ultimo_errore = f"Err: {str(e)[:20]}"
            continue
            
    return None, "OpenRouter", f"Esauriti. Ultimo err: {ultimo_errore}"

def genera_domande(blocco, cid):
    err_log = ""
    qs, engine, err = genera_gemini(blocco, cid)
    if qs is not None: return qs, engine, None
    if err: err_log += f"[GEM: {err}] "

    qs, engine, err = genera_groq(blocco, cid)
    if qs is not None: return qs, engine, None
    if err: err_log += f"[GROQ: {err}] "

    qs, engine, err = genera_openrouter(blocco, cid)
    if qs is not None: return qs, engine, None
    if err: err_log += f"[OR: {err}]"

    return [], "FALLITO", err_log.strip()

# =========================
# MAIN
# =========================

def main():
    os.system("cls" if os.name == "nt" else "clear")
    print_banner()
    
    template_scelto = scegli_grafica()
    print(f"\n\033[92m✅ Hai scelto il tema: {template_scelto}\033[0m\n")
    time.sleep(1)
    
    input_dir, output_dir = "input_docs", "output"
    os.makedirs(output_dir, exist_ok=True)
    
    db_path = os.path.join(output_dir, "database.json")
    path_mancati = os.path.join(output_dir, "blocchi_mancati.json")
    
    db = {"questions": [], "chapters": []}
    if os.path.exists(db_path):
        db = json.load(open(db_path, "r", encoding="utf-8"))

    usa_mancati = False
    mancati_data = []
    
    if os.path.exists(path_mancati):
        try:
            mancati_data = json.load(open(path_mancati, "r", encoding="utf-8"))
            if mancati_data and len(mancati_data) > 0:
                print(f"\033[93m⚠️  ATTENZIONE: Trovati {len(mancati_data)} blocchi salvati in 'blocchi_mancati.json' che avevano fallito.\033[0m")
                scelta = input("   Vuoi generare ORA le domande per questi blocchi mancanti? (s/n): ").strip().lower()
                if scelta == 's':
                    usa_mancati = True
                print("") 
        except: pass

    file_map = []
    total_blocks = 0

    if usa_mancati:
        print("\033[90m⚙️  Preparazione blocchi mancati in corso...\033[0m")
        grouped = {}
        for m in mancati_data:
            k = (m["f"], m["cid"])
            grouped.setdefault(k, []).append(m["testo"])
        
        for (f_nome, cid_val), blocchi_list in grouped.items():
            file_map.append((f_nome, cid_val, blocchi_list))
            total_blocks += len(blocchi_list)
            
    else:
        files = [f for f in os.listdir(input_dir) if f.endswith((".pdf",".docx",".pptx"))]
        if not files:
            print("\033[91m❌ Nessun file trovato nella cartella 'input_docs'\033[0m")
            return

        print("\033[90m⚙️  Controllo file e calcolo blocchi logici...\033[0m")
        capitoli_esistenti = {c["id"] for c in db.get("chapters", [])}

        for f in files:
            cid = f.split(".")[0].lower().replace(" ", "_")
            if cid in capitoli_esistenti: continue
            testo = estrai_testo(os.path.join(input_dir, f))
            blocchi = split_logico(testo)
            file_map.append((f, cid, blocchi))
            total_blocks += len(blocchi)

    if not file_map:
        if usa_mancati:
            print("\n\033[92m✅ Nessun blocco mancante da processare.\033[0m")
        else:
            print("\n\033[92m✅ Tutti i documenti sono già stati processati. Database aggiornato!\033[0m")
            
        print(f"\033[90m⚙️  Rigenero il sito web con il tema {template_scelto}...\033[0m\n")
        try:
            palette = random.choice(PALETTE_OPZIONI)
            with open(f"template/{template_scelto}", "r", encoding="utf-8") as t:
                html = t.read().replace("/* {{THEME_VARIABLES}} */", f":root {{ --ground: {palette['ground']}; --signal: {palette['signal']}; }}")
            with open(f"{output_dir}/index.html", "w", encoding="utf-8") as o:
                o.write(html)
        except Exception as e: print(f"Errore rigenerazione template: {e}")
        return

    print(f"\033[92m✓ Trovati {len(file_map)} file da elaborare ({total_blocks} blocchi totali).\033[0m\n")
    print("Inizio elaborazione. Non chiudere la finestra...\n")
    
    seen = set(q.get("q","") for q in db["questions"])
    processed = 0
    total_q = 0
    start_time = time.time()
    time_history = []
    api_stats = {"Gemini": 0, "Groq": 0, "OpenRouter": 0, "FALLITO": 0}
    is_first_draw = True
    ultimo_errore = ""
    
    nuovi_mancati = []

    for file_idx, (f, cid, blocchi) in enumerate(file_map):
        if not any(c.get("id") == cid for c in db.get("chapters", [])):
            db["chapters"].append({"id": cid, "nome_originale": f})

        tot_blocchi_file = len(blocchi)
        domande_nel_file = 0

        for i, b in enumerate(blocchi):
            b_start = time.time()
            blocco_idx = i + 1
            
            avg = sum(time_history[-10:]) / len(time_history[-10:]) if time_history else INTERVALLO
            eta_sec = (total_blocks - processed) * avg
            minuti, secondi = divmod(int(eta_sec), 60)
            eta_str = f"{minuti:02d}m {secondi:02d}s"
            
            elapsed_min = (time.time() - start_time) / 60.0
            speed = int(processed / elapsed_min) if elapsed_min > 0 else 0

            draw_dashboard(f, blocco_idx, tot_blocchi_file, domande_nel_file, "Pensiero...", eta_str, file_idx + 1, len(file_map), processed, total_blocks, total_q, speed, api_stats["Gemini"], api_stats["Groq"], api_stats["OpenRouter"], ultimo_errore, is_first_draw)
            is_first_draw = False

            qs, engine, current_error = genera_domande(b, cid)
            
            time_history.append(time.time() - b_start)
            processed += 1
            num_domande = len(qs)
            
            if current_error:
                ultimo_errore = current_error
            else:
                ultimo_errore = "" 
            
            if num_domande == 0:
                nuovi_mancati.append({"f": f, "cid": cid, "testo": b})
            
            domande_nel_file += num_domande
            total_q += num_domande
            
            if engine in api_stats: api_stats[engine] += 1

            elapsed_min = (time.time() - start_time) / 60.0
            speed = int(processed / elapsed_min) if elapsed_min > 0 else 0
            eta_sec = (total_blocks - processed) * avg
            minuti, secondi = divmod(int(eta_sec), 60)
            eta_str = f"{minuti:02d}m {secondi:02d}s"
            
            draw_dashboard(f, blocco_idx, tot_blocchi_file, domande_nel_file, engine, eta_str, file_idx + 1, len(file_map), processed, total_blocks, total_q, speed, api_stats["Gemini"], api_stats["Groq"], api_stats["OpenRouter"], ultimo_errore, False)

            for q in qs:
                if q["q"] not in seen:
                    db["questions"].append(q)
                    seen.add(q["q"])

    if usa_mancati:
        if len(nuovi_mancati) == 0 and os.path.exists(path_mancati):
            os.remove(path_mancati)
            print(f"\n\033[92m🗑️  Tutti i blocchi recuperati! Il file 'blocchi_mancati.json' è stato eliminato.\033[0m")
        elif len(nuovi_mancati) > 0:
            json.dump(nuovi_mancati, open(path_mancati, "w", encoding="utf-8"), indent=2)
    else:
        if len(nuovi_mancati) > 0:
            mancati_data.extend(nuovi_mancati)
            json.dump(mancati_data, open(path_mancati, "w", encoding="utf-8"), indent=2)

    json.dump(db, open(db_path,"w",encoding="utf-8"), indent=2)
    palette = random.choice(PALETTE_OPZIONI)
    
    with open(f"{output_dir}/data.js", "w", encoding="utf-8") as f:
        f.write(f"const BANK = {json.dumps(db['questions'], indent=2)};\n")
        meta = "const META = [\n"
        for i, c in enumerate(db["chapters"]):
            col = palette["cap_tones"][i % len(palette["cap_tones"])]
            nome_pulito = c.get("nome_originale", c["id"]).replace(".pdf", "").replace(".docx", "").replace(".pptx", "")
            meta += f'  {{ch:"{c["id"]}", num:"{i+1:02d}", nm:"{nome_pulito}", sb:"\", tone:"{col}"}},\n'
        f.write(meta + "];\n")

   # ... dopo aver caricato il template ...
    try:
        with open(template_path, "r", encoding="utf-8") as t:
            html = t.read()
        
        # INVECE di usare percorsi complessi, assicuriamoci che il CSS sia caricato correttamente
        # Se il tuo file è in output/index.html, il css deve stare in output/css/
        html = html.replace("/* {{THEME_VARIABLES}} */", f":root {{ --ground: {palette['ground']}; --signal: {palette['signal']}; }}")
        
        # FORZA il percorso CSS relativo corretto
        html = html.replace('href="css/', 'href="./css/') 
        
        with open(output_path, "w", encoding="utf-8") as o:
            o.write(html)
    except Exception as e:
        print(f"Errore: {e}")

    minuti_totali = int((time.time() - start_time) // 60)
    print_final_box(minuti_totali, len(file_map), total_blocks, total_q, len(db['questions']))
    
    if len(nuovi_mancati) > 0:
        print(f"\033[93m⚠️  Nota: {len(nuovi_mancati)} blocchi hanno generato 0 domande e sono stati salvati in 'blocchi_mancati.json'.\033[0m\n")

if __name__ == "__main__":
    main()